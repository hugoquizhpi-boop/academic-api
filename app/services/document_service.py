"""
Servicio de documentos: genera archivos .docx, .pdf y .pptx desde texto plano.
"""

import os
import uuid
from app.models.document_models import DocumentRequest

# Directorios de salida
MEDIA_DOCS_DIR = "media/documents"


def _generate_word(text: str, title: str, output_path: str):
    """Genera un archivo .docx usando python-docx."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Título del documento
    heading = doc.add_heading(title, level=1)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph()  # Espacio

    # Contenido: dividir por párrafos (líneas en blanco)
    paragraphs = text.strip().split("\n\n")
    for para_text in paragraphs:
        if para_text.strip():
            para = doc.add_paragraph(para_text.strip())
            para.style.font.size = Pt(11)

    doc.save(output_path)


def _generate_pdf(text: str, title: str, output_path: str):
    """Genera un archivo .pdf usando reportlab."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib import colors

    doc = SimpleDocTemplate(
        output_path,
        pagesize=A4,
        rightMargin=2.5*cm,
        leftMargin=2.5*cm,
        topMargin=2.5*cm,
        bottomMargin=2.5*cm
    )

    styles = getSampleStyleSheet()

    # Estilo para el título
    title_style = ParagraphStyle(
        "CustomTitle",
        parent=styles["Heading1"],
        fontSize=18,
        textColor=colors.HexColor("#1E3A5F"),
        spaceAfter=20,
        alignment=1  # Centrado
    )

    # Estilo para el cuerpo
    body_style = ParagraphStyle(
        "CustomBody",
        parent=styles["Normal"],
        fontSize=11,
        leading=16,  # Interlineado
        spaceAfter=12
    )

    story = []
    story.append(Paragraph(title, title_style))
    story.append(Spacer(1, 0.5*cm))

    # Dividir texto en párrafos
    paragraphs = text.strip().split("\n\n")
    for para_text in paragraphs:
        if para_text.strip():
            # Escapar caracteres especiales para reportlab
            safe_text = para_text.strip().replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe_text, body_style))

    doc.build(story)


def _generate_ppt(text: str, title: str, output_path: str):
    """
    Genera un archivo .pptx usando python-pptx.
    Divide el texto en slides automáticamente (1 párrafo = 1 slide).
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()

    # ── Slide de título ───────────────────────────────────────────────────────
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Documento generado automáticamente"

    # ── Slides de contenido ───────────────────────────────────────────────────
    content_layout = prs.slide_layouts[1]  # Título + contenido
    paragraphs = text.strip().split("\n\n")

    for i, para_text in enumerate(paragraphs):
        if not para_text.strip():
            continue

        slide = prs.slides.add_slide(content_layout)

        # Título del slide: primera línea del párrafo o "Sección N"
        lines = para_text.strip().split("\n")
        slide_title = lines[0][:50] if lines[0] else f"Sección {i+1}"
        slide_content = "\n".join(lines[1:]) if len(lines) > 1 else para_text.strip()

        slide.shapes.title.text = slide_title
        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True
        tf.text = slide_content

    prs.save(output_path)


async def generate_document(request: DocumentRequest) -> tuple[str, str]:
    """
    Genera el documento según el formato solicitado.

    Returns:
        Tupla (file_path, media_type) para enviarlo como respuesta de descarga
    """
    os.makedirs(MEDIA_DOCS_DIR, exist_ok=True)

    unique_id = uuid.uuid4().hex[:8]

    if request.format == "word":
        filename = f"document_{unique_id}.docx"
        output_path = os.path.join(MEDIA_DOCS_DIR, filename)
        _generate_word(request.text, request.title, output_path)
        media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

    elif request.format == "pdf":
        filename = f"document_{unique_id}.pdf"
        output_path = os.path.join(MEDIA_DOCS_DIR, filename)
        _generate_pdf(request.text, request.title, output_path)
        media_type = "application/pdf"

    elif request.format == "ppt":
        filename = f"document_{unique_id}.pptx"
        output_path = os.path.join(MEDIA_DOCS_DIR, filename)
        _generate_ppt(request.text, request.title, output_path)
        media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    else:
        raise ValueError(f"Formato no soportado: {request.format}")

    return output_path, media_type
